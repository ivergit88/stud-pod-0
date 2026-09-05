# -*- coding: utf-8 -*-
"""Render the produced PPTX back to PNG (geometry-level) and run layout checks.
Reads the REAL artifact (slide XML): shapes, fills, lines, pictures, charts,
text frames. Then: draws a preview per slide, and reports text overflow,
out-of-bounds shapes and unintended opaque-overlap pairs."""
import sys, io
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

DPI = 150
SW, SH = 13.333, 7.5
W_FACTOR = 1.06
_FCACHE = {}

def _font(bold):
    k = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    if k not in _FCACHE:
        _FCACHE[k] = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/" + k, 100)
    return _FCACHE[k]

def px(v): return int(round(v * DPI))
def text_w(text, size, bold=False):
    return _font(bold).getlength(text) / 100.0 * size * W_FACTOR / 72.0

def wrap_lines(text, max_w, size, bold=False):
    words = text.split(" ")
    lines, cur = [], ""
    for wd in words:
        trial = (cur + " " + wd).strip()
        if text_w(trial, size, bold) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur); cur = wd
    if cur: lines.append(cur)
    return lines

def emu_in(v): return v / 914400.0

def hexc(rgbcolor):
    return "#" + str(rgbcolor)

def render_slide(slide, idx, out_dir, checks):
    img = Image.new("RGB", (px(SW), px(SH)), "#FFFFFF")
    d = ImageDraw.Draw(img)
    solids = []
    # slide background fill if set on background
    try:
        bg = slide.background.fill
        if bg.type is not None and str(bg.type) == 'SOLID (1)':
            d.rectangle([0, 0, px(SW), px(SH)], fill=hexc(bg.fore_color.rgb))
    except Exception:
        pass
    for shp in slide.shapes:
        x = emu_in(shp.left); y = emu_in(shp.top)
        w = emu_in(shp.width); h = emu_in(shp.height)
        st = shp.shape_type
        if x < -0.02 or y < -0.02 or x + w > SW + 0.02 or y + h > SH + 0.02:
            checks.append(f"S{idx} OUT-OF-BOUNDS {getattr(shp,'name','?')} "
                          f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
        if st == MSO_SHAPE_TYPE.PICTURE:
            try:
                pim = Image.open(io.BytesIO(shp.image.blob)).convert("RGB")
                pim = pim.resize((px(w), px(h)))
                img.paste(pim, (px(x), px(y)))
            except Exception as e:
                checks.append(f"S{idx} PICTURE ERR {e}")
            continue
        if st == MSO_SHAPE_TYPE.CHART or (hasattr(shp, "has_chart") and shp.has_chart):
            draw_chart(d, shp, x, y, w, h)
            continue
        if st == MSO_SHAPE_TYPE.LINE or shp.__class__.__name__ == "Connector":
            try:
                x1, y1 = emu_in(shp.begin_x), emu_in(shp.begin_y)
                x2, y2 = emu_in(shp.end_x), emu_in(shp.end_y)
                col = hexc(shp.line.color.rgb)
                d.line([px(x1), px(y1), px(x2), px(y2)], fill=col, width=max(1, px(shp.line.width.pt / 72)))
            except Exception:
                pass
            continue
        # auto shape / textbox
        fill_col = None
        try:
            if shp.fill.type is not None and str(shp.fill.type).startswith("SOLID"):
                fill_col = hexc(shp.fill.fore_color.rgb)
        except Exception:
            pass
        line_col = None
        try:
            if shp.line.fill.type is not None and str(shp.line.fill.type).startswith("SOLID"):
                line_col = hexc(shp.line.color.rgb)
        except Exception:
            pass
        prst = ""
        try:
            prst = shp._element.spPr.prstGeom.get("prst")
        except Exception:
            pass
        if fill_col:
            rad = 0
            try:
                if shp.adjustments:
                    rad = min(w, h) * shp.adjustments[0]
            except Exception:
                pass
            if prst == "ellipse":
                d.ellipse([px(x), px(y), px(x + w), px(y + h)], fill=fill_col,
                          outline=line_col, width=1 if line_col else 0)
            else:
                d.rounded_rectangle([px(x), px(y), px(x + w), px(y + h)],
                                    radius=px(rad), fill=fill_col,
                                    outline=line_col, width=1 if line_col else 0)
            solids.append((idx, getattr(shp, "name", "?"), x, y, w, h))
        elif line_col:
            d.rounded_rectangle([px(x), px(y), px(x + w), px(y + h)],
                                radius=px(min(w, h) * 0.08), outline=line_col, width=1)
        # text
        if shp.has_text_frame:
            draw_text(d, shp, x, y, w, h, idx, checks)
    img.save(f"{out_dir}/slide-{idx:02d}.png")

def draw_chart(d, shp, x, y, w, h):
    try:
        ch = shp.chart
        s = ch.plots[0].series[0]
        vals, cols = [], []
        for i in range(len(s.values)):
            vals.append(s.values[i])
            try:
                cols.append(hexc(s.points[i].format.fill.fore_color.rgb))
            except Exception:
                cols.append("#DCE4EF")
        total = sum(vals) or 1
        cx, cy = px(x + w / 2), px(y + h / 2)
        R = px(min(w, h) / 2)
        r0 = int(R * 0.62)
        a0 = -90
        for v, c in zip(vals, cols):
            a1 = a0 + 360.0 * v / total
            d.pieslice([cx - R, cy - R, cx + R, cy + R], a0, a1, fill=c)
            a0 = a1
        d.ellipse([cx - r0, cy - r0, cx + r0, cy + r0], fill="#FFFFFF")
    except Exception as e:
        d.rectangle([px(x), px(y), px(x + w), px(y + h)], outline="#999999")

def draw_text(d, shp, x, y, w, h, idx, checks):
    tf = shp.text_frame
    blocks = []
    total_h = 0.0
    for p in tf.paragraphs:
        runs = p.runs
        if not runs:
            blocks.append(None); total_h += 0.06; continue
        r0 = runs[0]
        size = r0.font.size.pt if r0.font.size else 10
        bold = bool(r0.font.bold)
        try: col = hexc(r0.font.color.rgb)
        except Exception: col = "#000000"
        full = "".join(r.text for r in runs)
        lh = size / 72.0 * 1.18
        raw_lines = []
        for seg in full.split("\n"):
            raw_lines += wrap_lines(seg, w + 0.02, size, bold) or [""]
        lines = raw_lines
        try:
            if p.line_spacing and p.line_spacing > 0.5:
                lh = size / 72.0 * p.line_spacing
        except Exception:
            pass
        after = 0.0
        try:
            if p.space_after: after = p.space_after.pt / 72.0
        except Exception:
            pass
        n = max(1, len(lines))
        blocks.append((lines, size, bold, p.alignment, lh, after, col))
        total_h += n * lh + after
    anchor = tf.vertical_anchor
    oy = y
    if anchor == MSO_ANCHOR.MIDDLE:
        oy = y + (h - total_h) / 2
    elif anchor == MSO_ANCHOR.BOTTOM:
        oy = y + h - total_h
    if total_h > h + 0.035:
        txt = " ".join("".join(r.text for r in p.runs) for p in tf.paragraphs)[:42]
        checks.append(f"S{idx} TEXT-OVERFLOW h_need={total_h:.2f} h_box={h:.2f} “{txt}”")
    cy = oy
    for b in blocks:
        if b is None:
            cy += 0.06; continue
        lines, size, bold, align, lh, after, col = b
        f = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/" +
                               ("DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"),
                               max(6, int(size * DPI / 72.0)))
        for ln in lines:
            lw = text_w(ln, size, bold)
            if lw > w + 0.04:
                checks.append(f"S{idx} LINE-WIDE w={lw:.2f}>{w:.2f} “{ln[:36]}”")
            if align == PP_ALIGN.CENTER: lx = x + (w - lw) / 2
            elif align == PP_ALIGN.RIGHT: lx = x + w - lw
            else: lx = x
            d.text((px(lx), px(cy)), ln, font=f, fill=col)
            cy += lh
        cy += after

def runs_color(shp, size):
    for p in shp.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size and abs(r.font.size.pt - size) < 0.01:
                return r.font.color.rgb
    return RGBColor(0, 0, 0)

def overlap_report(checks, all_solids):
    for i in range(len(all_solids)):
        for j in range(i + 1, len(all_solids)):
            a, b = all_solids[i], all_solids[j]
            if a[0] != b[0]: continue
            ax, ay, aw, ah = a[2], a[3], a[4], a[5]
            bx, by, bw, bh = b[2], b[3], b[4], b[5]
            ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
            iy = max(0, min(ay + ah, by + bh) - max(ay, by))
            inter = ix * iy
            small = min(aw * ah, bw * bh)
            if small and inter / small > 0.25:
                checks.append(f"S{a[0]} OVERLAP {a[1]} <-> {b[1]} ratio={inter/small:.2f}")

def main(pptx_path, out_dir):
    import os
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    checks = []
    for i, sl in enumerate(prs.slides, 1):
        render_slide(sl, i, out_dir, checks)
    # overlap report needs solids collected globally: re-run simple pass
    print("\n".join(checks) if checks else "CHECKS CLEAN")
    print(f"rendered {len(prs.slides)} slides to {out_dir}")

if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else "render")

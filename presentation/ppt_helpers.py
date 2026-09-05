# -*- coding: utf-8 -*-
"""Helpers for building the Студенческий подряд deck (python-pptx) + a geometry
renderer/validator that reads the produced PPTX back and draws it with PIL."""
import io, math, re
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

SW, SH = 13.333, 7.5

# ------------------------------------------------------------------ palette
INK      = "0C1628"
BLUE     = "2864E8"
GREEN    = "10A77A"
PURPLE   = "7840EE"
CYAN     = "18B8D8"
ORANGE   = "E78A16"
GRAY     = "637089"
GRAYBLU  = "526984"
CARDBLU  = "EDF3FC"
CARDBLU2 = "EAF1FF"
CARDGRN  = "E8F7F2"
CARDPUR  = "F0EAFE"
CARDORG  = "FDF3E3"
BORDER   = "DCE4EF"
BG       = "F7F9FC"
WHITE    = "FFFFFF"
DARK     = "07182C"

FONT = "Arial"

def rgb(h): return RGBColor.from_string(h)

def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ------------------------------------------------------------------ shapes
def _set_radius(shape, r):
    try:
        shape.adjustments[0] = r
    except Exception:
        pass

def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None,
         dash=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            sp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if radius is not None:
        _set_radius(sp, radius)
    sp.shadow.inherit = False
    return sp

def line(slide, x1, y1, x2, y2, color=GRAY, w=1.0, dash=False):
    cnv = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cnv.line.color.rgb = rgb(color); cnv.line.width = Pt(w)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        cnv.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return cnv

# ------------------------------------------------------------------ text
def add_text(slide, x, y, w, h, paras, anchor="top", wrap=True):
    """paras: list of dicts(text,size,color,bold,align,after,before,italic,ls)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for i, p in enumerate(paras):
        run_p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                           "r": PP_ALIGN.RIGHT}[p.get("align", "l")]
        if p.get("after") is not None: run_p.space_after = Pt(p["after"])
        if p.get("before") is not None: run_p.space_before = Pt(p["before"])
        if p.get("ls") is not None:
            run_p.line_spacing = p["ls"]
        r = run_p.add_run()
        r.text = p["text"]
        f = r.font
        f.size = Pt(p.get("size", 10))
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.name = FONT
        f.color.rgb = rgb(p.get("color", INK))
    return tb

def T(text, size, color=INK, bold=False, align="l", after=None, before=None,
      italic=False, ls=None):
    return dict(text=text, size=size, color=color, bold=bold, align=align,
                after=after, before=before, italic=italic, ls=ls)

# ------------------------------------------------------------------ metrics
_FCACHE = {}
def _font(bold):
    k = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    if k not in _FCACHE:
        _FCACHE[k] = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/" + k, 100)
    return _FCACHE[k]

# DejaVu is ~4-6% wider than Arial for these scripts -> conservative estimate
W_FACTOR = 1.06

def text_w(text, size, bold=False):
    f = _font(bold)
    return f.getlength(text) / 100.0 * size * W_FACTOR / 72.0

def wrap_lines(text, max_w, size, bold=False):
    words = text.split(" ")
    lines, cur = [], ""
    for wd in words:
        trial = (cur + " " + wd).strip()
        if text_w(trial, size, bold) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur); cur = wd
    if cur: lines.append(cur)
    return lines

# ------------------------------------------------------------------ nav+footer
NAV = ["ЛОГИКА", "РЕЗУЛЬТАТИВНОСТЬ", "ИДЕЯ", "НОВИЗНА + ОБОСНОВАННОСТЬ",
       "ЛИЧНЫЙ ВКЛАД"]

def nav(slide, active):
    """active: index, or -1 for equal weight (slide 8 handled separately)."""
    y = 0.26; h = 0.36
    sizes = [text_w(s, 8, True) + 0.04 for s in NAV]
    arrow_w = text_w("→", 8, False) + 0.18
    pad = 0.16
    total = sum(sizes) + (len(NAV) - 1) * arrow_w + sum(
        pad * 2 if i == active else 0 for i in range(len(NAV)))
    x = (SW - total) / 2
    for i, s in enumerate(NAV):
        act = i == active
        w = sizes[i] + (pad * 2 if act else 0)
        if act:
            rect(slide, x, y - 0.06, w, h, fill=CARDBLU2, radius=0.35)
            rect(slide, x + 0.1, y + h - 0.1, w - 0.2, 0.045, fill=BLUE, radius=0.5)
        add_text(slide, x + (pad if act else 0), y, w - (pad * 2 if act else 0), h,
                 [T(s, 8, BLUE if act else GRAYBLU, True, "c")], anchor="middle")
        x += w
        if i < len(NAV) - 1:
            add_text(slide, x + 0.05, y, arrow_w - 0.1, h,
                     [T("→", 8, BORDER, False, "c")], anchor="middle")
            x += arrow_w

def footer(slide, source, num, dark=False, prefix="Источник: "):
    sc = "8FA3BC" if dark else GRAY
    add_text(slide, 0.75, 7.13, 10.6, 0.3,
             [T(prefix + source, 7, sc, False)])
    add_text(slide, 12.0, 7.1, 0.62, 0.3,
             [T(num, 8.5, sc, True, "r")])

def picture(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                    Inches(w) if w else None,
                                    Inches(h) if h else None)

# ------------------------------------------------------------------ donut
def donut(slide, x, y, size, value, color, rest=BORDER):
    cd = CategoryChartData()
    cd.categories = ["v", "r"]
    cd.add_series("s", (value, 100.0 - value))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y),
                                Inches(size), Inches(size), cd)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    plot = ch.plots[0]
    plot.has_data_labels = False
    try: plot.hole_size = 62
    except Exception: pass
    try: plot.first_slice_ang = 270
    except Exception: pass
    s0 = ch.plots[0].series[0]
    p0 = s0.points[0]; p0.format.fill.solid(); p0.format.fill.fore_color.rgb = rgb(color)
    p0.format.line.fill.background()
    p1 = s0.points[1]; p1.format.fill.solid(); p1.format.fill.fore_color.rgb = rgb(rest)
    p1.format.line.fill.background()
    # strip chart border / background
    try:
        ch.chart_style = 2
        el = gf.chart._chartSpace
        for tag in ("c:plotArea",):
            pass
    except Exception: pass
    try:
        from pptx.oxml.ns import qn as _qn
        spPr = gf.chart._chartSpace.get_or_add_chart().get_or_add_plotArea()
    except Exception: pass
    return gf

# ------------------------------------------------------------------ QR
def make_qr(text, path):
    import qrcode
    img = qrcode.make(text, box_size=10, border=1)
    img.save(path)
    return path

# ------------------------------------------------------------------ placeholder
def photo_placeholder(slide, x, y, w, h, label, fname, dark=False):
    rect(slide, x, y, w, h, fill=("12263C" if dark else CARDBLU),
         line=("2A4763" if dark else BORDER), dash=True, radius=0.04)
    add_text(slide, x, y + h / 2 - 0.34, w, 0.3,
             [T("ФОТО", 8, ("7FA6CE" if dark else GRAYBLU), True, "c")], anchor="middle")
    add_text(slide, x + 0.05, y + h / 2 - 0.05, w - 0.1, 0.26,
             [T(label, 7, ("7FA6CE" if dark else GRAYBLU), False, "c")], anchor="middle")
    add_text(slide, x + 0.05, y + h / 2 + 0.16, w - 0.1, 0.22,
             [T(fname, 6.2, ("5E82AC" if dark else GRAY), False, "c")], anchor="middle")
